"""
Generate the Planet B-612 pitch deck as a .pptx file.
Mirrors slides.html (the HTML deck) so the same 7 slides can be played on stage
from PowerPoint or Keynote.  Hackathon rules require .ppt/.keynote only.

Run inside the project's venv:
    source .venv-slides/bin/activate
    python build_slides.py
Produces:  PlanetB612_pitch.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ── Palette (matches WorldQuest app) ────────────────────────────────────────
BG          = RGBColor(0x07, 0x09, 0x1a)   # deep space
BG2         = RGBColor(0x0f, 0x11, 0x33)   # secondary panel
CARD        = RGBColor(0x12, 0x16, 0x32)   # glass card body
BORDER      = RGBColor(0x35, 0x3b, 0x55)
ACCENT      = RGBColor(0x00, 0xd4, 0xff)   # primary cyan
ACCENT2     = RGBColor(0x7c, 0x3a, 0xed)   # secondary violet
GOLD        = RGBColor(0xff, 0xd1, 0x66)
GREEN       = RGBColor(0x00, 0xff, 0x88)
TXT         = RGBColor(0xe6, 0xe9, 0xf5)
TXT2        = RGBColor(0x9a, 0xa0, 0xc0)
WHITE       = RGBColor(0xff, 0xff, 0xff)
BLACK       = RGBColor(0x00, 0x00, 0x00)

# ── Slide dimensions: 16:9 widescreen ──────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

# ── Helpers ─────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill, line=None, transparency=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.line.fill.background() if line is None else None
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    return s

def add_rounded(slide, x, y, w, h, fill, border=None, radius=0.05):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = radius
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s

def add_text(slide, x, y, w, h, text, size=18, color=TXT, bold=False,
             italic=False, align=PP_ALIGN.LEFT, font='Inter', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    return tb, p, r

def add_multi_text(slide, x, y, w, h, runs, size=18, color=TXT, bold=False,
                   align=PP_ALIGN.LEFT, font='Inter'):
    """runs = [(text, kwargs_dict)] — each can override color/bold/italic/size/font."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    for i, (text, kw) in enumerate(runs):
        if i == 0:
            r = p.add_run()
        else:
            r = p.add_run()
        r.text = text
        r.font.name = kw.get('font', font)
        r.font.size = Pt(kw.get('size', size))
        r.font.color.rgb = kw.get('color', color)
        r.font.bold = kw.get('bold', bold)
        r.font.italic = kw.get('italic', False)
    return tb

def fill_bg(slide, color=BG):
    """Set whole slide background to a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_glow_blob(slide, x, y, w, h, color):
    """Soft elliptical color blob to imitate cosmic gradient (semi-transparent)."""
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    # Make it semi-transparent via XML
    sp = s.fill.fore_color._xFill
    # Insert alpha (0..100000)
    alpha = etree.SubElement(sp, qn('a:alpha'))
    alpha.set('val', '14000')   # ~14% opacity
    return s

def add_stars(slide, seed=0):
    """Scatter a few small bright dots to imitate stars."""
    import random
    rng = random.Random(seed)
    for _ in range(60):
        cx = Emu(int(rng.uniform(0, SLIDE_W)))
        cy = Emu(int(rng.uniform(0, SLIDE_H)))
        sz = Emu(int(rng.uniform(10000, 30000)))
        d  = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, sz, sz)
        d.fill.solid()
        d.fill.fore_color.rgb = WHITE if rng.random() > 0.4 else ACCENT
        d.line.fill.background()
        # Light alpha
        sp = d.fill.fore_color._xFill
        alpha = etree.SubElement(sp, qn('a:alpha'))
        alpha.set('val', str(int(rng.uniform(10000, 50000))))

def add_topbar(slide, page_label):
    """Brand + page label at the very top."""
    add_text(slide, Inches(0.7), Inches(0.3), Inches(3), Inches(0.25),
             '● PLANET B-612', size=10, color=ACCENT, bold=True, font='JetBrains Mono')
    add_text(slide, Inches(9.6), Inches(0.3), Inches(3.1), Inches(0.25),
             page_label.upper(), size=10, color=TXT2, bold=False,
             font='JetBrains Mono', align=PP_ALIGN.RIGHT)

def add_eyebrow(slide, x, y, text):
    add_text(slide, x, y, Inches(6), Inches(0.3), text,
             size=11, color=ACCENT, bold=True, font='Inter')

def add_h2(slide, x, y, w, plain_left, accent_mid='', plain_right=''):
    tb = slide.shapes.add_textbox(x, y, w, Inches(1.2))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    for txt, col, bold in [(plain_left, TXT, True), (accent_mid, ACCENT, True), (plain_right, TXT, True)]:
        if not txt:
            continue
        r = p.add_run()
        r.text = txt
        r.font.name = 'Inter'
        r.font.size = Pt(36)
        r.font.color.rgb = col
        r.font.bold = bold
    return tb

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — HERO
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
fill_bg(s)
add_glow_blob(s, Inches(7), Inches(-2), Inches(8), Inches(8), ACCENT2)
add_glow_blob(s, Inches(-2), Inches(4), Inches(8), Inches(7), ACCENT)
add_stars(s, seed=1)
add_topbar(s, 'SuperAI Singapore · June 2026')

# Hero content
add_text(s, Inches(0.95), Inches(1.6), Inches(4), Inches(1.4),
         '🌍🦊', size=92, color=WHITE, font='Apple Color Emoji')

add_text(s, Inches(0.95), Inches(3.0), Inches(10), Inches(0.4),
         'A LITTLE-PRINCE-INSPIRED COMPANION APP',
         size=12, color=ACCENT, bold=True, font='JetBrains Mono')

add_text(s, Inches(0.95), Inches(3.45), Inches(10), Inches(1.4),
         'Planet B-612', size=78, color=ACCENT, bold=True)

add_multi_text(s, Inches(0.95), Inches(5.0), Inches(11), Inches(1.4), [
    ('A tiny planet that ', {'size': 22, 'color': TXT}),
    ('grows with your travels', {'size': 22, 'color': ACCENT, 'bold': True}),
    (', tended by a fox who waits for the stories you bring back.', {'size': 22, 'color': TXT}),
])

add_text(s, Inches(0.95), Inches(6.2), Inches(11), Inches(0.4),
         'VISION + AGENTS + 3D   ·   BUILT IN 1 DAY   ·   SINGAPORE-THEMED',
         size=11, color=TXT2, font='JetBrains Mono', bold=True)

# Chip pill at the bottom (live demo URL)
chip = add_rounded(s, Inches(0.95), Inches(6.65), Inches(4.6), Inches(0.4),
                   RGBColor(0x00, 0x35, 0x4a), border=ACCENT, radius=0.5)
add_text(s, Inches(0.95), Inches(6.65), Inches(4.6), Inches(0.4),
         '✦  ai-planet-b612-superai.vercel.app', size=12, color=ACCENT, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
         font='JetBrains Mono')

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
fill_bg(s)
add_glow_blob(s, Inches(8), Inches(-3), Inches(7), Inches(7), ACCENT)
add_stars(s, seed=2)
add_topbar(s, 'The Problem')

add_eyebrow(s, Inches(0.7), Inches(0.85), 'WHY THIS')
add_h2(s, Inches(0.7), Inches(1.25), Inches(12),
       'Travel memories ', 'deserve a home', "")
add_h2(s, Inches(0.7), Inches(1.85), Inches(12),
       "that isn't a camera roll.", '', '')

add_text(s, Inches(0.7), Inches(2.85), Inches(11.9), Inches(0.9),
         "We snap thousands of photos. Most are never seen again. We tell incredible stories — once — to a friend, then forget the detail. "
         "The trip ends, the world dissolves back into a folder named 2026/SG.",
         size=15, color=TXT2)

# Two cards
def card_at(x, title_emoji, title, items, accent=False):
    c = add_rounded(s, x, Inches(4.05), Inches(5.85), Inches(2.95), CARD,
                    border=ACCENT if accent else BORDER, radius=0.04)
    add_text(s, x + Inches(0.35), Inches(4.25), Inches(5.5), Inches(0.4),
             f'{title_emoji}  {title}', size=18, color=ACCENT, bold=True)
    tb = s.shapes.add_textbox(x + Inches(0.35), Inches(4.75), Inches(5.45), Inches(2.2))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = '▸  ' + it
        r.font.name = 'Inter'
        r.font.size = Pt(13)
        r.font.color.rgb = TXT

card_at(Inches(0.7), '📸', 'Today', [
    'Photos pile up in chronological folders',
    'No structure: sunset, noodles, temple wall — all equal pixels',
    'Cultural context lives only in the moment',
    'Travel ends → memories fade',
])
card_at(Inches(6.8), '🌍', 'Planet B-612', [
    'Each photo becomes a living 3D model on your planet',
    'An AI fox identifies it and tells you the deep local context',
    'Codex grows; planet evolves; fox writes diary entries',
    'Years later, you can still walk through the world you brought back',
], accent=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — CONCEPT (the loop)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
fill_bg(s)
add_glow_blob(s, Inches(-2), Inches(-3), Inches(8), Inches(8), ACCENT2)
add_stars(s, seed=3)
add_topbar(s, 'The Concept')

add_eyebrow(s, Inches(0.7), Inches(0.85), 'THE LOOP')
add_h2(s, Inches(0.7), Inches(1.25), Inches(12),
       'From ', 'snap', ' to ')
add_h2(s, Inches(0.7), Inches(1.85), Inches(12),
       '', 'grow', ' in four steps.')

# Quote (slim accent border on left)
add_rect(s, Inches(0.7), Inches(2.95), Emu(38100), Inches(1.4), ACCENT)  # 3pt = ~38100 emu
add_text(s, Inches(0.95), Inches(2.95), Inches(11.5), Inches(0.9),
         '"It is the time you have wasted for your rose that makes your rose so important."',
         size=20, color=TXT, italic=True)
add_text(s, Inches(0.95), Inches(3.85), Inches(11.5), Inches(0.4),
         '— Antoine de Saint-Exupéry, The Little Prince',
         size=11, color=TXT2, font='JetBrains Mono')

# Loop circles
loop_y = Inches(4.95)
circle_d = Inches(1.55)
spacing_x = Inches(2.85)
start_x = Inches(0.85)
steps = [('📷', 'CAPTURE', 'Snap anything\nthat catches your eye'),
         ('🤖', 'IDENTIFY', 'Qwen-VL +\nlocal context'),
         ('🌍', 'PLANT', 'Low-poly model\nspawns on planet'),
         ('🦊', 'FOX LIVES', 'Walks, writes diary,\nplans next wish')]

for i, (em, lbl, sub) in enumerate(steps):
    cx = start_x + spacing_x * i
    # Circle (ellipse w==h)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, loop_y, circle_d, circle_d)
    c.fill.solid()
    c.fill.fore_color.rgb = RGBColor(0x0c, 0x1f, 0x35)
    c.line.color.rgb = ACCENT
    c.line.width = Pt(1.5)
    # Emoji centered in circle
    add_text(s, cx, loop_y + Inches(0.15), circle_d, Inches(0.7),
             em, size=34, color=WHITE, align=PP_ALIGN.CENTER,
             font='Apple Color Emoji', anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, cx, loop_y + Inches(0.9), circle_d, Inches(0.3),
             lbl, size=10, color=ACCENT, bold=True, align=PP_ALIGN.CENTER,
             font='JetBrains Mono')
    # Sub-caption under circle
    add_text(s, cx - Inches(0.2), loop_y + circle_d + Inches(0.05),
             circle_d + Inches(0.4), Inches(0.7), sub,
             size=10, color=TXT2, align=PP_ALIGN.CENTER)
    # Arrow between steps
    if i < len(steps) - 1:
        ax = cx + circle_d + Inches(0.15)
        add_text(s, ax, loop_y + Inches(0.55), Inches(0.5), Inches(0.4),
                 '→', size=22, color=ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — EXPERIENCE (4 UI screens)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
fill_bg(s)
add_glow_blob(s, Inches(9), Inches(2), Inches(7), Inches(7), ACCENT)
add_stars(s, seed=4)
add_topbar(s, 'The Experience')

add_eyebrow(s, Inches(0.7), Inches(0.85), 'FOUR SURFACES, ONE WORLD')
add_h2(s, Inches(0.7), Inches(1.25), Inches(12),
       'Tend a planet. ', 'Read a journal.', ' Share a card.')

screens = [
    ('🌍', '3D Planet',
     'A trackball-rotatable low-poly sphere. Every photo adds a unique 3D model — building, tree, fox, fish — anchored to a meaningful surface spot. Hover to preview, drag to rearrange.',
     ['Three.js', '20 model templates', 'Non-overlap placement']),
    ('📷', 'Capture Modal',
     'One FAB at center-bottom — the only input. Pick a photo + 3 quick context chips (time / weather / company). A robust fallback lets you manually categorize if Vision fails.',
     ['Qwen-VL-Max', 'Fuzzy de-dupe', 'Cancel any time']),
    ('🌱', 'Codex Gallery',
     '20 category tabs (Landmarks, Buildings, Sacred Sites, Cuisine, Wildlife…). Click a card → highlights its 3D twin on the planet. Hover the model → floating tooltip the other way.',
     ['Tabbed Pokédex', 'Newest first', 'Cross-link']),
    ('🦊', 'Fox Journal & Share',
     'The fox writes Tabikaeru-style diary entries about your finds and its own little planet life. One-tap "Share Today" exports a 1080×1350 card with planet ID, level, and the day’s discoveries.',
     ['Tabikaeru loop', 'Canvas export', 'Planet ID']),
]
positions = [
    (Inches(0.7), Inches(2.8)),
    (Inches(6.9), Inches(2.8)),
    (Inches(0.7), Inches(5.0)),
    (Inches(6.9), Inches(5.0)),
]
SCREEN_W, SCREEN_H = Inches(5.75), Inches(2.05)

for (em, title, desc, tags), (x, y) in zip(screens, positions):
    add_rounded(s, x, y, SCREEN_W, SCREEN_H, CARD, border=BORDER, radius=0.05)
    add_text(s, x + Inches(0.3), y + Inches(0.18), Inches(0.6), Inches(0.5),
             em, size=22, color=WHITE, font='Apple Color Emoji', anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.85), y + Inches(0.18), Inches(4.5), Inches(0.5),
             title, size=15, color=ACCENT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.3), y + Inches(0.75), SCREEN_W - Inches(0.55), Inches(1.0),
             desc, size=10.5, color=TXT)
    # Tag chips
    tag_y = y + SCREEN_H - Inches(0.4)
    tx = x + Inches(0.3)
    for t in tags:
        chip_w = Inches(0.05 + 0.085 * len(t))
        chip = add_rounded(s, tx, tag_y, chip_w, Inches(0.27),
                           RGBColor(0x06, 0x2a, 0x3a), border=ACCENT, radius=0.5)
        add_text(s, tx, tag_y, chip_w, Inches(0.27), t,
                 size=8.5, color=ACCENT, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tx += chip_w + Inches(0.08)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ARCHITECTURE (hierarchy + flow)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
fill_bg(s)
add_glow_blob(s, Inches(-3), Inches(2), Inches(8), Inches(7), ACCENT2)
add_stars(s, seed=5)
add_topbar(s, 'Architecture')

add_eyebrow(s, Inches(0.7), Inches(0.85), 'AGENT HIERARCHY & FLOW')
add_h2(s, Inches(0.7), Inches(1.25), Inches(12), 'One router. ', 'Two pipelines.', ' One autonomous fox.')
add_text(s, Inches(0.7), Inches(2.4), Inches(8.6), Inches(0.7),
         "A single Orchestrator sits between the user and every specialist. "
         "It dispatches each input into the right pipeline — identification or planning — "
         "while Fox Life runs in parallel as its own autonomous loop.",
         size=12, color=TXT2)

# ─── Reusable agent-node drawer ─────────────────────────────────────────
def arch_node(x, y, w, h, em, name, role, sub_label='',
              card_fill=RGBColor(0x14, 0x1b, 0x3d), border=BORDER,
              name_color=ACCENT, dot_color=ACCENT, name_size=12.5,
              role_size=8.5):
    add_rounded(s, x, y, w, h, card_fill, border=border, radius=0.06)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, x + w - Inches(0.28), y + Inches(0.16),
                              Inches(0.12), Inches(0.12))
    dot.fill.solid(); dot.fill.fore_color.rgb = dot_color; dot.line.fill.background()
    add_text(s, x + Inches(0.22), y + Inches(0.12), Inches(0.6), Inches(0.4),
             em, size=18, color=WHITE, font='Apple Color Emoji', anchor=MSO_ANCHOR.MIDDLE)
    # Name line (with optional sub-label suffix)
    if sub_label:
        add_multi_text(s, x + Inches(0.22), y + Inches(0.5), w - Inches(0.45), Inches(0.3),
                       [(name, {'color': name_color, 'bold': True, 'size': name_size}),
                        (' ' + sub_label, {'color': TXT2, 'bold': False, 'size': 8.5})])
    else:
        add_text(s, x + Inches(0.22), y + Inches(0.5), w - Inches(0.45), Inches(0.3),
                 name, size=name_size, color=name_color, bold=True)
    add_text(s, x + Inches(0.22), y + Inches(0.85), w - Inches(0.4), h - Inches(0.95),
             role, size=role_size, color=TXT2)

# ─── Layout coordinates ─────────────────────────────────────────────────
# Main flow area: x = 0.7 → 9.4 (3 cols), with side panel at x = 9.6 → 12.8
COL_W = Inches(2.8)
COL_GAP = Inches(0.15)
COL_X = [Inches(0.7),
         Inches(0.7) + COL_W + COL_GAP,
         Inches(0.7) + 2 * (COL_W + COL_GAP)]

# Entry node — User (tier 0), small pill centered above orchestrator (col 1)
entry_w = Inches(2.4)
entry_x = COL_X[1] + (COL_W - entry_w) / 2
add_rounded(s, entry_x, Inches(3.2), entry_w, Inches(0.5),
            RGBColor(0x0a, 0x10, 0x28), border=BORDER, radius=0.5)
add_text(s, entry_x, Inches(3.2), entry_w, Inches(0.5),
         '👤  User · text / photo / tap', size=11, color=TXT, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font='Inter')

# Down arrow from user to orchestrator
add_text(s, entry_x, Inches(3.75), entry_w, Inches(0.3),
         '▼', size=14, color=ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# TIER 1 · Orchestrator (col 1, highlighted "core")
arch_node(COL_X[1], Inches(4.05), COL_W, Inches(1.05),
          '🧠', 'Orchestrator',
          'Classifies intent (message / photo / task / landmark) and dispatches to one pipeline. Owns shared planet context.',
          card_fill=RGBColor(0x06, 0x33, 0x4a),
          border=ACCENT,
          name_size=13)

# Branching arrows from orchestrator (down-left to Vision, down to Story, down-right to Planner)
# Three small downward arrows to indicate branching
for ax_inch in [COL_X[0] + COL_W / 2, COL_X[1] + COL_W / 2, COL_X[2] + COL_W / 2]:
    add_text(s, ax_inch - Inches(0.2), Inches(5.15), Inches(0.4), Inches(0.3),
             '▼', size=13, color=ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Tier-2 specialists row
TIER2_Y = Inches(5.45)
TIER2_H = Inches(1.15)
arch_node(COL_X[0], TIER2_Y, COL_W, TIER2_H,
          '👁️', 'Vision', 'Qwen-VL-Max → 1 of 20 categories + structured details. Receives existing codex for de-dupe.',
          sub_label='· photo pipeline')
arch_node(COL_X[1], TIER2_Y, COL_W, TIER2_H,
          '📖', 'Story', "Called by Orchestrator after Vision succeeds. Generates deep cultural / historical Insight in the fox's voice.",
          sub_label='· enriches Vision')
arch_node(COL_X[2], TIER2_Y, COL_W, TIER2_H,
          '✈️', 'Planner', 'SG local-guide brain. Builds 30-90 min quests grounded in time, weather, festivals — asks Sponsor.',
          sub_label='· text pipeline')

# Sponsor sub-call from Planner (below Planner col 3)
SPONSOR_Y = TIER2_Y + TIER2_H + Inches(0.15)
add_text(s, COL_X[2] + COL_W / 2 - Inches(0.2), SPONSOR_Y - Inches(0.18),
         Inches(0.4), Inches(0.2), '▼', size=11, color=GOLD,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
arch_node(COL_X[2], SPONSOR_Y, COL_W, Inches(0.85),
          '🎁', 'Sponsor',
          'Planner injects clearly-badged sponsored quests when keywords match (SuperAI / Tiger / Mr. Coconut).',
          sub_label='· sub-call',
          card_fill=RGBColor(0x2a, 0x21, 0x0a),
          border=GOLD,
          name_color=GOLD,
          dot_color=GOLD,
          name_size=12)

# ─── Side panel · Fox Life autonomous loop ──────────────────────────────
side_x = Inches(9.85)
side_w = Inches(3.05)
side_y = Inches(2.4)
side_h = Inches(4.3)
add_rounded(s, side_x, side_y, side_w, side_h,
            RGBColor(0x1f, 0x18, 0x09), border=GOLD, radius=0.04)
# "AUTONOMOUS LOOP" cap
add_rounded(s, side_x + Inches(0.4), side_y - Inches(0.15), Inches(1.55), Inches(0.3),
            BG, border=GOLD, radius=0.5)
add_text(s, side_x + Inches(0.4), side_y - Inches(0.15), Inches(1.55), Inches(0.3),
         'AUTONOMOUS LOOP', size=8, color=GOLD, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font='JetBrains Mono')

add_text(s, side_x + Inches(0.25), side_y + Inches(0.25), Inches(0.6), Inches(0.5),
         '🦊', size=24, color=WHITE, font='Apple Color Emoji', anchor=MSO_ANCHOR.MIDDLE)
add_text(s, side_x + Inches(0.25), side_y + Inches(0.75), side_w - Inches(0.5), Inches(0.4),
         'Fox Life', size=15, color=GOLD, bold=True)
add_text(s, side_x + Inches(0.25), side_y + Inches(1.15), side_w - Inches(0.5), Inches(1.1),
         "Runs in its own timer-driven loop, independent of the Orchestrator. "
         "Picks activities based on bond level + current planet state.",
         size=10, color=TXT2)

# Inner flow box
inner_y = side_y + Inches(2.4)
add_rounded(s, side_x + Inches(0.25), inner_y, side_w - Inches(0.5), Inches(1.7),
            RGBColor(0x05, 0x06, 0x18), border=BORDER, radius=0.05)
flow_lines = [
    'every ~15 min  →',
    'pick activity  →',
    'write 60–110 word diary  →',
    'sometimes return with keepsake  →',
    'write back to planet context',
]
flow_tb = s.shapes.add_textbox(side_x + Inches(0.4), inner_y + Inches(0.15),
                                side_w - Inches(0.7), Inches(1.45))
ftf = flow_tb.text_frame
ftf.margin_left = ftf.margin_right = Emu(0)
ftf.margin_top = ftf.margin_bottom = Emu(0)
ftf.word_wrap = True
for i, line in enumerate(flow_lines):
    p = ftf.paragraphs[0] if i == 0 else ftf.add_paragraph()
    p.space_after = Pt(2)
    r = p.add_run()
    r.text = line
    r.font.name = 'JetBrains Mono'
    r.font.size = Pt(9)
    r.font.color.rgb = GOLD if i == 0 else TXT2
    r.font.bold = (i == 0)

# Tier labels (left-edge faint markers)
def tier_label(y, text):
    add_text(s, Inches(0.7), y, Inches(2.5), Inches(0.22), text,
             size=8.5, color=TXT2, bold=True, font='JetBrains Mono')

tier_label(Inches(3.0), '— TIER 0 · ENTRY —')
tier_label(Inches(3.95), '— TIER 1 · ROUTER —')
tier_label(Inches(5.32), '— TIER 2 · SPECIALISTS —')

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — TECH + FEATURES
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
fill_bg(s)
add_glow_blob(s, Inches(9), Inches(-3), Inches(7), Inches(8), ACCENT)
add_glow_blob(s, Inches(-2), Inches(4), Inches(7), Inches(7), ACCENT2)
add_stars(s, seed=6)
add_topbar(s, 'Problems & Solutions')

add_eyebrow(s, Inches(0.7), Inches(0.85), 'PROBLEMS WE HIT · HOW WE FIXED THEM')
add_h2(s, Inches(0.7), Inches(1.25), Inches(12), 'Six hard problems. ', 'Six shipped fixes.', '')

# Each entry: (emoji, title, problem, solution)
probs = [
    ('🐛', 'Laksa → "landmark"',
     'Vision mis-classified hawker food as buildings, then UI silently showed fake fallback data.',
     'Negative-constraint prompt with food examples; on failure show "AI couldn\'t read this" + Retry / Manual buttons — never fake.'),
    ('⏱️', '2-minute API hangs',
     'Server logs showed Qwen calls aborting at ~30 ms — wrong req.on(\'close\') handler killed every request.',
     'Swapped to res.on(\'close\'), added 50s server timeout + 60s client timeout + true Cancel button.'),
    ('🌍', 'Models stacking on each other',
     'New items spawned on top of existing landmarks — a pile, not a planet.',
     'Spherical spiral search using great-circle distance picks the nearest free spot — guaranteed non-overlap.'),
    ('🧭', 'Buildings tilted on the sphere',
     'Models leaned at random Y-angles after placement; flagpoles aimed anywhere.',
     'Build orthonormal basis per spot (radial-up / world-north / east) → quaternion. Upright everywhere.'),
    ('♻️', '"Kaya Toast" ≠ "咖椰吐司"',
     'AI returned slightly different names for the same dish → duplicates in the codex.',
     'Pass existing codex into Vision prompt + fuzzy normalize (strip punctuation, prefixes) before dedupe.'),
    ('✨', 'Highlight hidden by modal',
     'Clicking a codex card highlighted the 3D model, but the full-screen detail modal covered the planet.',
     'Right-docked the modal, lowered overlay opacity, added halo + floating arrow mesh on top.'),
]
GRID_X, GRID_Y = Inches(0.7), Inches(2.4)
COL_W, ROW_H = Inches(3.95), Inches(1.85)
COL_GAP, ROW_GAP = Inches(0.2), Inches(0.18)

for i, (em, title, problem, solution) in enumerate(probs):
    col, row = i % 3, i // 3
    x = GRID_X + col * (COL_W + COL_GAP)
    y = GRID_Y + row * (ROW_H + ROW_GAP)
    add_rounded(s, x, y, COL_W, ROW_H, CARD, border=BORDER, radius=0.05)
    add_text(s, x + Inches(0.25), y + Inches(0.12), Inches(0.5), Inches(0.4),
             em, size=16, color=WHITE, font='Apple Color Emoji', anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.78), y + Inches(0.12), COL_W - Inches(0.95), Inches(0.4),
             title, size=12.5, color=ACCENT, bold=True, anchor=MSO_ANCHOR.MIDDLE)

    # Problem line
    add_multi_text(s, x + Inches(0.25), y + Inches(0.6), COL_W - Inches(0.45), Inches(0.55),
                   [('Problem  ', {'color': RGBColor(0xff, 0x6b, 0x8a), 'bold': True, 'size': 8.5, 'font': 'JetBrains Mono'}),
                    (problem, {'color': TXT, 'size': 9})])
    # Solution line
    add_multi_text(s, x + Inches(0.25), y + Inches(1.15), COL_W - Inches(0.45), Inches(0.6),
                   [('Fix  ', {'color': GREEN, 'bold': True, 'size': 8.5, 'font': 'JetBrains Mono'}),
                    (solution, {'color': TXT, 'size': 9})])

# Tech stack chips at the bottom
chip_data = [
    ('Three.js r128', ACCENT),
    ('Qwen-VL-Max  ·  vision', ACCENT),
    ('Qwen-Max  ·  planner / story / fox', ACCENT),
    ('Vanilla JS  ·  no framework', TXT2),
    ('Express proxy', TXT2),
    ('localStorage  ·  zero infra', TXT2),
]
chip_y = Inches(6.5)
cx = Inches(0.7)
for txt, col in chip_data:
    cw = Inches(0.4 + 0.10 * len(txt))
    add_rounded(s, cx, chip_y, cw, Inches(0.4),
                RGBColor(0x0e, 0x14, 0x32), border=BORDER, radius=0.5)
    add_text(s, cx, chip_y, cw, Inches(0.4), txt, size=10, color=col,
             bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             font='JetBrains Mono')
    cx += cw + Inches(0.1)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ROADMAP / CLOSING
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
fill_bg(s)
add_glow_blob(s, Inches(2), Inches(-2), Inches(9), Inches(9), ACCENT2)
add_glow_blob(s, Inches(8), Inches(4), Inches(7), Inches(7), ACCENT)
add_stars(s, seed=7)
add_topbar(s, "What's Next")

add_eyebrow(s, Inches(0.7), Inches(0.85), 'ROADMAP')
add_h2(s, Inches(0.7), Inches(1.25), Inches(12),
       'From ', 'Singapore-day-one', ' to your every trip.')

roadmap = [
    ('1', 'Multi-city',
     'Lift the Singapore-only knowledge to a generic local-guide template. Activate per detected city: Tokyo, Bangkok, Bali, Lisbon.'),
    ('2', 'Multiplayer planets',
     'Share your planet via planet ID. Friends can "visit" — leave footprints, swap species, even gift a model from their travels.'),
    ('3', 'Real-world rewards',
     'Sponsored quests verified by photo → unlock real coupons. Already wired with SuperAI 2026, Tiger, Mr. Coconut. Scale via brand partnerships.'),
]
GRID_X, GRID_Y = Inches(0.7), Inches(3.0)
COL_W, ROW_H = Inches(4.0), Inches(2.4)
COL_GAP = Inches(0.2)
for i, (num, title, desc) in enumerate(roadmap):
    x = GRID_X + i * (COL_W + COL_GAP)
    y = GRID_Y
    add_rounded(s, x, y, COL_W, ROW_H, CARD, border=BORDER, radius=0.04)
    # Number badge (circle) — top-left, overlapping
    badge_d = Inches(0.7)
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.4), y - Inches(0.3),
                                badge_d, badge_d)
    badge.fill.solid(); badge.fill.fore_color.rgb = ACCENT
    badge.line.color.rgb = BG
    badge.line.width = Pt(2)
    add_text(s, x + Inches(0.4), y - Inches(0.3), badge_d, badge_d,
             num, size=20, color=BLACK, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, x + Inches(0.4), y + Inches(0.6), COL_W - Inches(0.8), Inches(0.5),
             title, size=18, color=ACCENT, bold=True)
    add_text(s, x + Inches(0.4), y + Inches(1.05), COL_W - Inches(0.8), Inches(1.3),
             desc, size=12.5, color=TXT)

# Closer
add_text(s, Inches(0.7), Inches(5.95), Inches(12), Inches(0.5),
         '🦊  Thank you', size=28, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.3),
         'Built in one day for SuperAI Singapore 2026',
         size=12, color=TXT2, italic=True, align=PP_ALIGN.CENTER, font='Inter')
add_multi_text(s, Inches(0.7), Inches(6.95), Inches(12), Inches(0.3), [
    ('🌐  Live demo  ',                              {'color': TXT2, 'size': 11, 'font': 'JetBrains Mono'}),
    ('ai-planet-b612-superai.vercel.app',            {'color': ACCENT, 'size': 11, 'bold': True, 'font': 'JetBrains Mono'}),
    ('     📦  Code  ',                              {'color': TXT2, 'size': 11, 'font': 'JetBrains Mono'}),
    ('github.com/zhsg2024-hub/AI-Planet-B612',       {'color': ACCENT, 'size': 11, 'bold': True, 'font': 'JetBrains Mono'}),
], align=PP_ALIGN.CENTER)

# ── Save ────────────────────────────────────────────────────────────────────
out = 'PlanetB612_pitch.pptx'
prs.save(out)
print(f'✓ Wrote {out}')
print(f'  Slides: {len(prs.slides)}')
print(f'  Aspect: 16:9 widescreen ({prs.slide_width.inches} x {prs.slide_height.inches} in)')
